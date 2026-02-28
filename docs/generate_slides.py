#!/usr/bin/env python3
# /// script
# requires-python = ">=3.13"
# dependencies = [
#     "python-pptx",
# ]
# ///
# mypy: disable-error-code="no-untyped-def, no-untyped-call"
"""Generate the Claude Code Workshop PowerPoint presentation."""

import os
from dataclasses import dataclass, field
from pptx import Presentation  # type: ignore[import-not-found]
from pptx.dml.color import RGBColor  # type: ignore[import-not-found]
from pptx.util import Pt  # type: ignore[import-not-found]

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.dirname(SCRIPT_DIR)
TEMPLATE = os.path.join(SCRIPT_DIR, "template-converted.pptx")
OUTPUT = os.path.join(SCRIPT_DIR, "workshop-slides.pptx")

# Styling constants
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
BG_IMAGE = os.path.join(REPO_ROOT, "images", "powerpoint_background.png")

# Template layout indices (from template-converted.pptx)
LAYOUT_TITLE = 0
LAYOUT_SECTION = 2
LAYOUT_TWO_COLUMN = 5
LAYOUT_CONTENT = 6
LAYOUT_CLOSING = 7

# Bullet type alias
Bullet = str | tuple[str, str, int] | dict[str, str | int]  # pyright: ignore[reportGeneralTypeIssues]


@dataclass
class SlideData:
    """Base class for all slide types."""

    title: str
    notes: str = ""


@dataclass
class TitleSlide(SlideData):
    subtitle: str = ""


@dataclass
class SectionSlide(SlideData):
    layout: int = -1  # Set after LAYOUT_SECTION is defined

    def __post_init__(self) -> None:
        if self.layout == -1:
            self.layout = LAYOUT_SECTION


@dataclass
class ContentSlide(SlideData):
    bullets: list[Bullet] = field(default_factory=list)


@dataclass
class ImageSlide(SlideData):
    image: str = ""  # Relative to REPO_ROOT


@dataclass
class TwoColumnSlide(SlideData):
    left: list[Bullet] = field(default_factory=list)
    right: list[Bullet] = field(default_factory=list)


# All 34 slides in exact order
SLIDES: list[SlideData] = [
    # Slide 1: Title slide
    TitleSlide(
        title="Claude Code in Depth",
        subtitle="Totally not a chat bot:\nharness your context for increased velocity",
        notes="Welcome to the Claude Code workshop. This session takes you from zero to a fully operational agentic delivery system in five progressive modules. Each module builds on the last — by the end, you'll have working orchestration commands, agent teams, and CI/CD integration. The focus is on principles of agentic coding, not just the project we build.",
    ),
    # Slide 2: Prerequisites
    ContentSlide(
        title="Prerequisites",
        bullets=[
            ("Git: ", "Pre-installed on macOS/Linux — Windows users must install", 0),
            ("uv: ", "Python package manager", 0),
            ("Claude Code CLI: ", "Anthropic's agentic coding tool", 0),
            (
                "AWS Profile: ",
                "SSO profile with Bedrock access (InvokeModel permission)",
                0,
            ),
        ],
        notes="These should all be installed before the workshop begins. If you need help you can reach out in the meeting chat or flag down an assistant. ",
    ),
    # Slide 3: Bedrock Configuration
    ContentSlide(
        title="Bedrock Configuration",
        bullets=[
            "In ~/.claude/settings.json, set env variables:",
            {"text": 'CLAUDE_CODE_USE_BEDROCK: "1"', "level": 1},
            {"text": 'AWS_REGION: "us-west-2"', "level": 1},
            {"text": 'AWS_PROFILE: "your-sso-profile.BedrockAccessRole"', "level": 1},
            "Pin model versions (prevents breaking changes):",
            {"text": "Sonnet: us.anthropic.claude-sonnet-4-6", "level": 1},
            {"text": "Haiku: us.anthropic.claude-haiku-4-5-20251001-v1:0", "level": 1},
            {"text": "Opus: us.anthropic.claude-opus-4-6-v1", "level": 1},
            'awsAuthRefresh: "aws sso login --profile=${AWS_PROFILE}"',
        ],
        notes="If you have used the configure.py this was done behind the scenes.",
    ),
    # Slide 4: Workshop Flow
    ContentSlide(
        title="What We'll Build",
        bullets=[
            "M1: Scaffold a Typer CLI app from a natural language prompt",
            "M2: Build a simple CLI LLM tool called todd",
            "M3: Build phase commands, agents, hooks, and skills",
            "M4: Build orchestration commands from PRD",
            "M5: Run parallel agent teams in isolated worktrees",
            "M6 (Take-home): Build todd into a Claude Code clone",
        ],
        notes="Frame this as a hands-on journey. Where we'll progressively build up capabilities and complexity. Transition: Let's start with what Claude Code actually is.",
    ),
    # Slide 5: How Modules Work
    ContentSlide(
        title="How Modules Work",
        bullets=[
            (
                "Each module has a guide: ",
                "`modules/moduleN.md` — follow it step by step",
                0,
            ),
            (
                "When you finish a module: ",
                "run `/module` in Claude Code to advance",
                0,
            ),
            (
                "What `/module` does: ",
                "merges your branch forward into the next module branch",
                0,
            ),
            (
                "Your work carries forward: ",
                "each module builds on everything you've built so far",
                0,
            ),
        ],
        notes="Orient participants before they start. The module files are the source of truth for exercises — they follow those, not the slides. The /module command is the transition mechanism: it merges the current branch into the next, so nothing is lost. Participants work at their own pace; the command handles the git logistics.",
    ),
    # Slide 6: Getting Started
    ContentSlide(
        title="Getting Started",
        bullets=[
            (
                "Open VS Code: ",
                "use the integrated terminal (`` Ctrl+` `` / `Cmd+J`) side by side with the code editor",
                0,
            ),
            (
                "Launch Claude: ",
                "run `claude` from the project root in the integrated terminal",
                0,
            ),
            (
                "Side-by-side workflow: ",
                "terminal split lets you see Claude output and code changes simultaneously",
                0,
            ),
            (
                "VS Code extension: ",
                "alternatively, install the Claude Code extension from the VS Code marketplace",
                0,
            ),
        ],
        notes="This is the 'how to begin' reference slide. Participants should already have VS Code installed (covered in SETUP.md). Walk them through opening the integrated terminal and launching Claude from the project root. The side-by-side workflow is critical — they need to see both Claude's output and the files it's editing. The extension is an alternative for participants who prefer a GUI integration over the terminal.",
    ),
    # Slide 7: What Is Claude Code?
    ContentSlide(
        title="What Is Claude Code?",
        bullets=[
            (
                "Agentic terminal tool: ",
                "Claude runs in your terminal with full filesystem access",
                0,
            ),
            (
                "Built-in tools: ",
                "Read, Write, Edit, Bash, Glob, Grep, Task (subagents)",
                0,
            ),
            (
                "MCP servers: ",
                "External tool integrations via Model Context Protocol",
                0,
            ),
            (
                "Extensibility layers: ",
                "commands, skills, hooks, agents — all defined in .claude/",
                0,
            ),
            (
                "Plugin marketplace: ",
                "Discover, install, and share packaged extensions via /plugin",
                0,
            ),
        ],
        notes="Claude Code is not a chatbot — it's an agentic tool that operates on your codebase. It reads files, writes code, runs commands, and manages its own workflow through subagents. The extensibility layers — MCP, commands, skills, hooks, agents — are what make it a platform, not just a tool. Plugins package these extensions for distribution: a single plugin can bundle commands, skills, agents, and hooks into one installable unit. Plugin marketplaces — both the public Anthropic marketplace and private enterprise ones — let teams share and discover these packages. Each extension point gets dedicated coverage in later modules. Transition: Before diving into the modules, let's establish the vocabulary.",
    ),
    # Slide 8: Quality Tenets — Foundations
    ContentSlide(
        title="Quality Tenets: Foundations",
        bullets=[
            ("1. Verify your work: ", "Tests and expected outputs are the single highest-leverage thing you can provide", 0),
            ("2. Be specific: ", "Reference files, constraints, patterns — vague prompts produce vague output", 0),
            ("3. CLAUDE.md + hooks: ", "Persistent memory plus automated gates — deterministic quality from commit one", 0),
            ("4. Context is finite: ", "Manage aggressively — performance degrades as the window fills", 0),
        ],
        notes=(
            "These four tenets are foundational — they apply from your very first Claude session. "
            "Each one maps to concepts you'll learn hands-on:\n\n"
            "Vocabulary preview: 'Context window' is the finite space for prompts, tools, and history. "
            "'Context rot' is quality degradation as this fills. 'Context poisoning' is when bad info "
            "enters and corrupts reasoning. These failure modes are why Tenet 4 matters.\n\n"
            "Module mapping: #1 starts in M1 (manual check) and deepens in M2 (automated test). "
            "#2 is demonstrated through PRD-driven workflows starting in M2. "
            "#3 is the M1 headline (CLAUDE.md + pre-commit) and deepens in M3 (PostToolUse hooks). "
            "#4 is the M2 headline (context visualization, subagent isolation).\n\n"
            "Transition: The next four tenets build on these foundations."
        ),
    ),
    # Slide 9: Quality Tenets — At Scale
    ContentSlide(
        title="Quality Tenets: At Scale",
        bullets=[
            ("5. Explore → Plan → Code: ", "Separate thinking from doing — read the codebase before writing to it", 0),
            ("6. Progressive disclosure: ", "Right context, right scope, right time — don't front-load everything", 0),
            ("7. Agent design: ", "Composable workers with model, tools, scope — agent design is API design", 0),
            ("8. Scale with isolation: ", "Worktrees, teams, headless, sandboxing — isolate to scale safely", 0),
        ],
        notes=(
            "These four tenets emerge as you move from individual use to team-scale orchestration.\n\n"
            "Vocabulary preview: 'Progressive disclosure' means surfacing information incrementally. "
            "'Back pressure' means the system actively resists bad output (hooks, TDD, plan mode). "
            "'Sandboxing' isolates agent operations to limit blast radius. "
            "'Delegation' offloads work to subagents for context savings.\n\n"
            "Module mapping: #5 is the M2 headline (plan mode workflow). "
            "#6 is the M3 headline (four-layer context system). "
            "#7 is the M3 headline (custom agents and custom subagents) extended in M4 (team workers). "
            "#8 is the M4 headline (worktrees + teams) extended in M5 (headless/CI).\n\n"
            "Transition: Let's start building. Module 1."
        ),
    ),
    # Slide 10: Module 1 Section Header
    SectionSlide(
        title="Module 1: Project Scaffolding",
        notes="Module 1 establishes the baseline: set up the tools, configure quality gates, and let Claude build a simple Typer CLI app from a natural language prompt. The workshop starts slow on purpose — scaffolding a CLI is boring but demonstrates Claude operating on your behalf. Claude is doing, not planning. Pre-commit hooks are back pressure (automated quality gates). CLAUDE.md shapes Claude's behavior across sessions. These patterns scale. Transition: What we do in this module.",
    ),
    # Slide 11: M1 Let's get to work
    ContentSlide(
        title="Let's get to work",
        bullets=[
            "Open your terminal, run `claude` from the project root, then open `modules/module1.md`",
            "Open `modules/module1.md` for our list of tasks",
        ],
        notes="Participants follow the steps in module1.md at their own pace. Flag down an assistant if you get stuck.\n\nPresenter talking points:\n- Install uv (Python package manager)\n- Configure Claude CLI (Bedrock profile)\n- Scaffold a Typer CLI app from a natural language prompt — Claude reads uv docs, decides structure, writes code, tests pass\n- Pre-commit hooks enforce quality gates automatically — ruff, mypy, bandit, etc.\n- Verify Claude followed your CLAUDE.md instructions by inspecting the generated code\n\nThe key insight: Claude is not a chatbot. It's an autonomous tool that can read, write, run commands, make decisions. Pre-commit hooks provide back pressure. CLAUDE.md provides persistent memory.\n\nTransition: Describe what you want, Claude handles the rest. Pre-commit hooks are back pressure before we even name it. CLAUDE.md shaped everything Claude just did.",
    ),
    # Slide 12: M1 Summary
    ContentSlide(
        title="Summarizing what we have seen",
        bullets=[
            ("Pre-commit hooks: ", "Automated quality gates enforce standards", 0),
            ("CLAUDE.md: ", "Persistent memory shapes behavior across sessions", 0),
            ("Simple prompt → full app: ", "Claude scaffolds, writes, tests autonomously", 0),
            ("Back pressure: ", "Hooks reject bad output before it commits", 0),
        ],
        notes="Describe what you want, Claude handles the rest. Pre-commit hooks are back pressure before we even name it. CLAUDE.md shaped everything Claude just did. Transition: Module 2 introduces context management and planning.",
    ),
    # Slide 13: Module 2 Section Header
    SectionSlide(
        title="Module 2: MCP, Plan Mode, and the Agent SDK",
        notes="Module 2 introduces the critical concept of context as a finite resource. You'll install your first MCP server, see how MCP tools load on demand via ToolSearch, switch models for planning, and build a real feature through plan mode. The key progression: Module 1 was about Claude doing things for you. Module 2 is about Claude thinking before doing. Transition: What we do in this module.",
    ),
    # Slide 14: M2 What We Do
    ContentSlide(
        title="Let's get to work",
        bullets=[
            "Open your terminal, run `claude` from the project root, then open `modules/module2.md`",
            "Open `modules/module2.md` for our list of tasks",
        ],
        notes="Participants follow the steps in module2.md at their own pace. Flag down an assistant if you get stuck.",
    ),
    # Slide 15: The Context Window (with image)
    ImageSlide(
        title="The Context Window",
        image="images/context_window_utilization.png",
        notes="This is the most important conceptual slide. Walk through each layer. System prompt and CLAUDE.md are always present — they're the 'tax' on every interaction. MCP tool definitions are deferred via ToolSearch — they load into context only when Claude discovers them. This is progressive disclosure at the tool level. Conversation history is the biggest consumer — this is why /compact and /clear matter. File contents are loaded on demand but can be huge. The 'available space' is what Claude has for actual reasoning. When it shrinks too much, that's context rot in action. Transition: Now let's see how subagents protect this window.",
    ),
    # Slide 16: Subagent Context Isolation (with image)
    ImageSlide(
        title="Subagent Context Isolation",
        image="images/subagent_context_savings.png",
        notes="This diagram shows how subagents protect the main context window. When Claude needs to search 40+ files or do deep research, it spawns a subagent — typically Haiku for cheap exploration, Opus for architecture decisions. Each subagent gets its own fresh context window. The search noise (file contents, failed matches, irrelevant results) stays in the subagent's window. Only a compact summary flows back to the main conversation. This is the primary defense against context rot: delegate the noisy work, keep the main window clean. Point out the three subagent types: two Explore agents running Haiku for fast, cheap searches, and one Plan agent running Opus for deeper reasoning. Transition: Module 3 builds the orchestration layer.",
    ),
    # Slide 17: M2 Summary
    ContentSlide(
        title="Summarizing what we have seen",
        bullets=[
            (
                "Finite context window: ",
                "Prompts, files, tools, history all share one space",
                0,
            ),
            (
                "MCP tools load on demand: ",
                "ToolSearch defers definitions until needed — no idle cost",
                0,
            ),
            (
                "Plan mode = read-only gate: ",
                "Think before writing, prevents costly reverts",
                0,
            ),
            (
                "Subagents isolate context: ",
                "Fresh windows prevent cross-contamination",
                0,
            ),
            ("Model selection: ", "Opus / Sonnet / Haiku for different task types", 0),
        ],
        notes="Context is finite and precious. Rot happens gradually (long sessions). Poisoning happens suddenly (one bad assumption). Subagents are the primary defense against both — each gets a fresh context window. MCP tools are deferred via ToolSearch — no idle context cost. Once loaded, definitions do consume context, so selective loading matters. This is progressive disclosure at the tool level. Model selection is about cost-quality tradeoffs, not just 'use the best model.' Transition: Module 3 builds the orchestration layer.\n\nTenet tracker: #4 (context is finite) and #5 (explore-plan-code) are the headlines. #1 (verify) gets its first test exercise. #2 (be specific) is demonstrated through the PRD workflow.",
    ),
    # Slide 18: Module 3 Section Header
    SectionSlide(
        title="Module 3: Extending Claude Code",
        notes="Module 3 builds the full extensibility layer. We create commands, skills, hooks, rules, and custom agents — all in .claude/. The key insight: all of this is markdown and configuration — no code. Claude's behavior is shaped by prose instructions, not programming. This is prompt-driven orchestration. Transition: What we do.",
    ),
    # Slide 19: M3 Let's get to work
    ContentSlide(
        title="Let's get to work",
        bullets=[
            "Open your terminal, run `claude` from the project root, then open `modules/module3.md`",
            "Open `modules/module3.md` for our list of tasks",
        ],
        notes="Participants follow the steps in module3.md at their own pace. Flag down an assistant if you get stuck.\n\nPresenter talking points:\n- Generate Architecture Reference Document with Mermaid diagrams — Claude as documentation tool, no code written, just analysis\n- Install document-skills plugin from the marketplace\n- Plan orchestration commands from PRD (docs/prds/adw-commands.md) — Claude reads .claude/ scaffolding to understand patterns, explores phase commands, skills, agents before planning\n- Build 4 orchestration commands from PRD: /feature (single-agent, 7 sequential phases), /bug (single-agent, 6 phases), /team:feature (multi-agent, parallel analysis + coordinated implementation), /team:bug (multi-agent, parallel analysis, 6 phases)\n- Create custom agents (.claude/agents/ + YAML frontmatter) — code-reviewer exercise teaches agent anatomy, then three ADW specialist agents\n- Explore PostToolUse lint hook — back pressure after every write",
    ),
    # Slide 20: The Four-Layer Context System (with image)
    ImageSlide(
        title="The Four-Layer Context System",
        image="images/fourlayer_context_system.png",
        notes="This is the architectural insight of Module 3. The four layers form a progressive disclosure system. CLAUDE.md is always present (broad but essential). Rules narrow to specific directories. Skills narrow to specific task types. Hooks narrow to specific actions. The key insight: context is delivered at the moment it's needed, at the right scope. This prevents context bloat while ensuring Claude always has what it needs. Transition: Let's dig deeper into when progressive disclosure works and when it doesn't.",
    ),
    # Slide 21: Progressive Disclosure (with image)
    ImageSlide(
        title="Progressive Disclosure",
        image="images/progressive_disclosure.png",
        notes="Progressive disclosure is the difference between a productive session and context rot. The four-layer funnel: CLAUDE.md (always-on) → rules/ (directory-scoped) → skills (task-scoped) → hooks (event-scoped). Each layer narrows scope and timing. Anti-patterns include: dumping 100KB+ in system prompt, loading all reference docs 'just in case', same context for every workflow phase, hiding context the agent needs now. Nielsen Norman research shows 3+ disclosure levels cause confusion even for humans. Progressive disclosure doesn't mean hiding information — it means delivering it at the right time. If an agent needs reference docs NOW, load them NOW. Transition: The other side of quality — back pressure.",
    ),
    # Slide 22: Back Pressure (with image)
    ImageSlide(
        title="Back Pressure",
        image="images/back_pressure_layers.png",
        notes="Back pressure is the system's resistance to bad output. Pre-commit hooks (M1) are the first layer — reject malformed commits. PostToolUse hooks (M3) are the second layer — inspect write operations before they land. Plan mode (M2) is architectural back pressure — think before writing. Tests (M1-M2) are verification back pressure — prove it works before you commit it. The goal: detect failure as early as possible. The later a failure is caught, the more expensive it is to fix. Hooks are cheap back pressure. They run synchronously, provide immediate feedback, and cost nothing if output is correct. Transition: Module 4 builds on this foundation by adding team-based orchestration.",
    ),
    # Slide 23: M3 Summary
    ContentSlide(
        title="Summarizing what we have seen",
        bullets=[
            (
                "Four-layer context system: ",
                "CLAUDE.md → rules/ → skills/ → hooks/",
                0,
            ),
            (
                "Progressive disclosure: ",
                "Deliver context at the right time and scope",
                0,
            ),
            (
                "Back pressure: ",
                "Hooks and gates reject bad output early",
                0,
            ),
            (
                "Custom agents: ",
                "YAML frontmatter + markdown instructions",
                0,
            ),
            (
                "Orchestration commands: ",
                "Phase-driven workflows via /command",
                0,
            ),
        ],
        notes="Module 3 is the turning point. We went from 'Claude does things' (M1) and 'Claude plans things' (M2) to 'Claude orchestrates workflows' (M3). The four-layer context system is the architectural win. Back pressure is the quality win. Custom agents are the extensibility win. Orchestration commands are the productivity win. Everything from here scales on this foundation.\n\nTenet tracker: #6 (progressive disclosure) is the M3 headline. #3 (CLAUDE.md + hooks) deepened with PostToolUse. #7 (agent design) introduced through custom agent YAML.",
    ),
    # Slide 24: Module 4 Section Header
    SectionSlide(
        title="Module 4: Building Orchestration Commands",
        notes="Module 4 builds the orchestration commands from a PRD. You'll use plan mode to architect four delivery workflows — /feature, /bug, /team:feature, /team:bug — and let Claude read the existing phase commands to understand the patterns before building. The key insight: Claude researches the scaffolding before planning, producing commands that correctly compose the primitives. This is prompt-driven orchestration. Transition: What we do.",
    ),
    # Slide 25: M4 Let's get to work
    ContentSlide(
        title="Let's get to work",
        bullets=[
            "Open your terminal, run `claude` from the project root, then open `modules/module4.md`",
            "Open `modules/module4.md` for our list of tasks",
        ],
        notes="Participants follow the steps in module4.md at their own pace. Flag down an assistant if you get stuck.\n\nPresenter talking points:\n- Activate plan mode, give Claude the ADW PRD\n- Claude reads existing phase commands to understand invocation patterns\n- Review the plan for all four orchestration commands\n- Approve and let Claude build: /feature, /bug, /team:feature, /team:bug\n- Each command composes the phase primitives in different ways",
    ),
    # Slide 26: M4 Summary
    ContentSlide(
        title="Summarizing what we have seen",
        bullets=[
            ("Orchestration commands: ", "Compose phase primitives into end-to-end workflows", 0),
            ("Plan mode as research: ", "Claude explores scaffolding before proposing a plan", 0),
            ("Four delivery variants: ", "/feature, /bug, /team:feature, /team:bug", 0),
            ("Prompt-driven composition: ", "Commands defined in markdown, no code required", 0),
        ],
        notes="Module 4 produced four orchestration commands by composing the seven phase primitives. Plan mode was used not for code but for configuration — Claude explored the codebase, understood the patterns, and produced a plan for markdown files. This is prompt-driven orchestration: no code written, just configuration. Transition: Module 5 runs these commands against real PRDs.",
    ),
    # Slide 27: Module 5 Section Header
    SectionSlide(
        title="Module 5: Team Orchestration",
        notes="Module 5 scales from single-agent workflows to multi-agent teams. You'll create worker agents that operate in parallel, coordinate via leader agents, and execute in isolated worktrees. The key insight: team orchestration is just phase-driven workflows with parallel execution. Each worker gets its own fresh worktree and context window. Leaders delegate, workers execute, results merge. This is how you scale Claude to multiple simultaneous work streams. Transition: What we do.",
    ),
    # Slide 28: M5 Let's get to work
    ContentSlide(
        title="Let's get to work",
        bullets=[
            "Open your terminal, run `claude` from the project root, then open `modules/module5.md`",
            "Open `modules/module5.md` for our list of tasks",
        ],
        notes="Participants follow the steps in module5.md at their own pace. "
              "Flag down an assistant if you get stuck.\n\n"
              "Presenter talking points:\n"
              "- Plan team worker agents from PRD (docs/prds/team-workers.md)\n"
              "- Create 3 worker agents: backend-dev, frontend-dev, test-engineer\n"
              "- Build orchestrator commands: /team:feature, /team:bug\n"
              "- Wire worktree isolation into team commands\n"
              "- Test parallel execution, observe worktree isolation\n"
              "- Inspect agent coordination: leader delegates, workers execute",
    ),
    # Slide 29: Worktree Isolation (with image)
    ImageSlide(
        title="Worktree Isolation",
        image="images/worktree_isolation.png",
        notes="This diagram shows how worktrees isolate parallel work. Each worker agent gets its own worktree — a separate working directory linked to the same git repository. Workers can write, commit, and test without interfering with each other. The leader agent operates in the main worktree and coordinates merges. This is the primary isolation mechanism for team orchestration. It prevents file conflicts, context contamination, and race conditions. Point out the three worktrees (backend, frontend, test) and the main worktree (leader). This is Module 5. Transition: Let's discuss managing Claude Code projects long-term.",
    ),
    # Slide 30: Managing Projects with Claude Code
    ContentSlide(
        title="Managing Projects with Claude Code",
        bullets=[
            ("CLAUDE.md audits: ", "Stale instructions actively mislead — review after major changes", 0),
            ("CI/CD with headless mode: ", "`claude -p /<command>` in pipelines for automated reviews and ADWs", 0),
            ("Automated delivery: ", "Orchestrators chain phase commands — same pattern as M4, in CI", 0),
            ("Cost & session hygiene: ", "`/cost`, `/stats`, `--max-budget-usd` to prevent surprise bills", 0),
            ("Checkpointing: ", "`Esc+Esc` / `/rewind` to recover from wrong directions", 0),
        ],
        notes="This is the high-level 'what comes next' discussion slide. CLAUDE.md audits "
              "tie back to the anti-patterns taught in M3 — now applied at maintenance scale. "
              "CI/CD utilization uses the same headless mode pattern participants saw in M5's "
              "orchestrators — `claude -p` with predefined commands automates ADWs or individual "
              "lifecycle phases (review, test, document) in GitHub Actions. Cost tracking and "
              "checkpointing are operational hygiene for sustained usage. Participants will use "
              "all of these skills in the M6 capstone project.",
    ),
    # Slide 31: The Progression (with image)
    ImageSlide(
        title="The Progression",
        image="images/progression.png",
        notes="This slide shows the full six-module progression. M1: CLAUDE.md + pre-commit (foundational quality). M2: Context management + subagents (efficiency). M3: Four-layer context + phase commands (scale). M4: Orchestration commands from PRD (composition). M5: Team workers + worktrees (parallelization). M6 (take-home): Build todd into a Claude Code clone. Each module builds on the last. The principles — progressive disclosure, back pressure, isolation, delegation — compound. By M5, you have a fully operational agentic delivery system. Transition: Let's review how the eight tenets mapped to this progression.",
    ),
    # Slide 32: Tenets — Recap
    ContentSlide(
        title="Tenets of Quality Output — Recap",
        bullets=[
            ("1. Verify your work: ", "tests and expected outputs", 0),
            ("2. Be specific: ", "reference files, constraints, patterns", 0),
            ("3. CLAUDE.md + hooks: ", "persistent memory + automated gates", 0),
            ("4. Context is finite: ", "performance degrades as it fills", 0),
            ("5. Explore → Plan → Code: ", "read before writing", 0),
            ("6. Progressive disclosure: ", "right context, right scope, right time", 0),
            ("7. Agent design: ", "composable workers with clear scope", 0),
            ("8. Scale with isolation: ", "worktrees, teams, sandboxing", 0),
        ],
        notes=(
            "This is the payoff slide. Walk through each tenet and ask participants to recall "
            "the specific exercises. Every tenet from slides 6-7 now has hands-on coverage.\n\n"
            "Tenets 1-4 (Foundations): verification, specificity, CLAUDE.md + hooks, context management.\n\n"
            "Tenets 5-8 (At Scale): plan mode discipline, progressive disclosure, agent design, isolation.\n\n"
            "Sources: code.claude.com/docs/en/best-practices, "
            "simonwillison.net (agentic engineering), "
            "github.com/anthropics/prompt-eng-interactive-tutorial"
        ),
    ),
    # Slide 33: Module 6 Homework
    ContentSlide(
        title="Homework: Build a Claude Code Clone",
        bullets=[
            "Module 6 is a self-paced capstone project",
            "Open `modules/module6.md` for 5 milestone PRDs",
            "Use your ADW tools — `/feature` and `/team:feature` work here",
        ],
        notes="Module 6 is take-home. Participants extend todd from a single-shot "
              "prompt forwarder into an interactive agentic tool — a miniature Claude "
              "Code clone. Five milestones: REPL, Tool Use, Streaming, CLAUDE.md "
              "Loading, Session Persistence. Each has a PRD in docs/prds/. Estimated "
              "2-4 hours across multiple sessions. Use `claude --resume` to pick up "
              "where you left off.",
    ),
    # Slide 34: Closing
    SectionSlide(
        title="Questions?",
        layout=LAYOUT_CLOSING,
        notes="Open the floor for questions. If time permits, offer to demo any specific concept live. Repo is available for participants to continue working through modules independently. Point participants to Module 6 — the capstone project where they build todd into a Claude Code clone using everything they've learned. Point them to code.claude.com for official documentation.",
    ),
]


def style_slide(slide, prs):
    """Set background image and dark text on all slide elements."""
    # Add background image covering full slide, sent to back
    pic = slide.shapes.add_picture(BG_IMAGE, 0, 0, prs.slide_width, prs.slide_height)
    # Move picture to back of z-order
    sp_tree = slide.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)

    # Set all text to dark color
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.color.rgb = DARK_TEXT
                for run in paragraph.runs:
                    run.font.color.rgb = DARK_TEXT


def fill_bullets(placeholder, bullets):
    """Fill a placeholder with bullet items."""
    tf = placeholder.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(bullet, tuple):
            bold_text, normal_text, level = bullet
            p.level = level
            run = p.add_run()
            run.text = bold_text
            run.font.bold = True
            if normal_text:
                run2 = p.add_run()
                run2.text = normal_text
        elif isinstance(bullet, dict):
            p.text = bullet["text"]
            p.level = bullet.get("level", 0)
        else:
            p.text = bullet
            p.level = 0


def delete_original_slides(prs, count):
    """Delete the first N slides (original template slides)."""
    # We must delete from the beginning, always removing index 0
    # because indices shift after each removal
    for _ in range(count):
        sldIdLst = prs.slides._sldIdLst
        rNs = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
        sldId_elem = sldIdLst[0]  # Always remove first slide
        rId = sldId_elem.get(rNs + "id")
        sldIdLst.remove(sldId_elem)
        prs.part.drop_rel(rId)


def set_notes_font_size(prs):
    """Set speaker notes font size to 18pt for readability."""
    for slide in prs.slides:
        if slide.has_notes_slide:
            for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                for run in paragraph.runs:
                    run.font.size = Pt(18)


def add_title_slide(prs, layout_idx, title_text, subtitle_text=None, notes_text=""):
    """Add a slide with centered title and optional subtitle."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            ph.text = title_text
        elif ph.placeholder_format.idx == 1 and subtitle_text:  # Subtitle
            ph.text = subtitle_text

    if notes_text:
        notes = slide.notes_slide
        notes.notes_text_frame.text = notes_text

    style_slide(slide, prs)
    return slide


def add_section_header(prs, layout_idx, title_text, notes_text=""):
    """Add a section header slide with centered title only."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            ph.text = title_text

    if notes_text:
        notes = slide.notes_slide
        notes.notes_text_frame.text = notes_text

    style_slide(slide, prs)
    return slide


def add_content_slide(prs, layout_idx, title_text, bullets, notes_text=""):
    """Add a slide with title and bullet list content."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    title_ph = None
    content_ph = None

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            title_ph = ph
        elif ph.placeholder_format.idx in (1, 10, 11):  # Content area
            content_ph = ph

    if title_ph:
        title_ph.text = title_text

    if content_ph:
        fill_bullets(content_ph, bullets)

    if notes_text:
        notes = slide.notes_slide
        notes.notes_text_frame.text = notes_text

    style_slide(slide, prs)
    return slide


def add_image_slide(prs, layout_idx, title_text, image_path, notes_text=""):
    """Add a slide with title and a full-width image in the content area."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    title_ph = None
    content_ph = None

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            title_ph = ph
        elif ph.placeholder_format.idx in (1, 10, 11):
            content_ph = ph

    if title_ph:
        title_ph.text = title_text

    # Remove the content placeholder and add the image in its place
    if content_ph:
        left = content_ph.left
        top = content_ph.top
        width = content_ph.width
        height = content_ph.height
        # Remove placeholder
        sp = content_ph._element
        sp.getparent().remove(sp)
        # Add image
        slide.shapes.add_picture(image_path, left, top, width, height)

    if notes_text:
        notes = slide.notes_slide
        notes.notes_text_frame.text = notes_text

    style_slide(slide, prs)
    return slide


def add_two_column_slide(
    prs, layout_idx, title_text, left_bullets, right_bullets, notes_text=""
):
    """Add a two-column slide."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    # Title
    if 0 in placeholders:
        placeholders[0].text = title_text

    # Left content (idx=1)
    if 1 in placeholders:
        fill_bullets(placeholders[1], left_bullets)

    # Right content (idx=2)
    if 2 in placeholders:
        fill_bullets(placeholders[2], right_bullets)

    if notes_text:
        notes = slide.notes_slide
        notes.notes_text_frame.text = notes_text

    style_slide(slide, prs)
    return slide


def render_title(prs, slide_data):
    add_title_slide(
        prs, LAYOUT_TITLE, slide_data.title, slide_data.subtitle, slide_data.notes
    )


def render_section(prs, slide_data):
    add_section_header(prs, slide_data.layout, slide_data.title, slide_data.notes)


def render_content(prs, slide_data):
    add_content_slide(
        prs, LAYOUT_CONTENT, slide_data.title, slide_data.bullets, slide_data.notes
    )


def render_image(prs, slide_data):
    image_path = os.path.join(REPO_ROOT, slide_data.image)
    add_image_slide(prs, LAYOUT_CONTENT, slide_data.title, image_path, slide_data.notes)


def render_two_column(prs, slide_data):
    add_two_column_slide(
        prs,
        LAYOUT_TWO_COLUMN,
        slide_data.title,
        slide_data.left,
        slide_data.right,
        slide_data.notes,
    )


RENDERERS = {
    TitleSlide: render_title,
    SectionSlide: render_section,
    ContentSlide: render_content,
    ImageSlide: render_image,
    TwoColumnSlide: render_two_column,
}


def main():
    prs = Presentation(TEMPLATE)

    # Print available layouts for debugging
    print("Available layouts:")
    for i, layout in enumerate(prs.slide_layouts):
        placeholders_info = [
            (ph.placeholder_format.idx, ph.name) for ph in layout.placeholders
        ]
        print(f"  Layout {i}: {layout.name} — placeholders: {placeholders_info}")
    print()

    original_slide_count = len(prs.slides)
    print(f"Original template has {original_slide_count} slides\n")

    for slide_data in SLIDES:
        renderer = RENDERERS[type(slide_data)]
        renderer(prs, slide_data)

    print(f"Deleting {original_slide_count} original template slides...")
    delete_original_slides(prs, original_slide_count)
    set_notes_font_size(prs)

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"\nSaved presentation with {len(prs.slides)} slides to {OUTPUT}")


if __name__ == "__main__":
    main()
